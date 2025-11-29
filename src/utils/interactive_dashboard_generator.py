"""
Interactive Dashboard Generator

Generates D3.js interactive dashboards conforming to DOCUMENTATION-FORMAT-SPEC-v1.0.

Author: Asif Hussain
License: Source-Available (Use Allowed, No Contributions)
"""

from pathlib import Path
from typing import Dict, List, Optional, Any
import json
from datetime import datetime
import subprocess
import sys


class InteractiveDashboardGenerator:
    """
    Reference implementation of D3.js dashboard generator.
    
    Features:
    - 5-tab structure (Overview, Visualizations, Diagrams, Data, Recommendations)
    - D3.js force-directed graphs
    - Chart.js time series
    - Mermaid diagram embedding
    - Export to PDF/PNG/PPTX
    """
    
    def __init__(self):
        """Initialize dashboard generator."""
        self.version = "1.0.0"
        self.template_path = Path(__file__).parent.parent.parent / "templates" / "interactive-dashboard-template.html"
        
    def generate_dashboard(
        self,
        title: str,
        data: Dict[str, Any],
        output_file: str
    ) -> bool:
        """
        Generate interactive D3.js dashboard HTML file.
        
        Args:
            title: Dashboard title
            data: Dashboard data conforming to format-validation-schema.json
            output_file: Output HTML file path
            
        Returns:
            bool: True if generation successful
            
        Example:
            >>> generator = InteractiveDashboardGenerator()
            >>> data = {
            ...     "metadata": {"generatedAt": "2025-11-28T14:30:00Z", ...},
            ...     "overview": {"executiveSummary": "...", ...},
            ...     "visualizations": {...},
            ...     "diagrams": [...],
            ...     "dataTable": [...],
            ...     "recommendations": [...]
            ... }
            >>> generator.generate_dashboard("My Dashboard", data, "output.html")
            True
        """
        try:
            # Validate data structure
            self._validate_data(data)
            
            # Load template
            template = self._load_template()
            
            # Inject data
            html = self._inject_data(template, title, data)
            
            # Write output
            output_path = Path(output_file)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_text(html, encoding="utf-8")
            
            return True
            
        except Exception as e:
            print(f"Dashboard generation failed: {e}")
            return False
    
    def _validate_data(self, data: Dict[str, Any]) -> None:
        """
        Validate data structure against schema.
        
        Args:
            data: Dashboard data to validate
            
        Raises:
            ValueError: If data structure invalid
        """
        required_keys = ["metadata", "overview", "visualizations", "diagrams", "dataTable", "recommendations"]
        
        for key in required_keys:
            if key not in data:
                raise ValueError(f"Missing required key: {key}")
        
        # Validate metadata
        metadata_keys = ["generatedAt", "version", "operationType", "author"]
        for key in metadata_keys:
            if key not in data["metadata"]:
                raise ValueError(f"Missing metadata.{key}")
        
        # Validate overview
        overview_keys = ["executiveSummary", "keyMetrics", "statusIndicator"]
        for key in overview_keys:
            if key not in data["overview"]:
                raise ValueError(f"Missing overview.{key}")
        
        # Validate key metrics count
        if not (5 <= len(data["overview"]["keyMetrics"]) <= 10):
            raise ValueError(f"keyMetrics must have 5-10 items, got {len(data['overview']['keyMetrics'])}")
        
        # Validate visualizations
        viz_keys = ["forceGraph", "timeSeries"]
        for key in viz_keys:
            if key not in data["visualizations"]:
                raise ValueError(f"Missing visualizations.{key}")
        
        # Validate force graph structure
        if "nodes" not in data["visualizations"]["forceGraph"]:
            raise ValueError("Missing forceGraph.nodes")
        if "links" not in data["visualizations"]["forceGraph"]:
            raise ValueError("Missing forceGraph.links")
        
        # Validate diagrams array
        if not isinstance(data["diagrams"], list) or len(data["diagrams"]) == 0:
            raise ValueError("diagrams must be non-empty array")
        
        # Validate data table array
        if not isinstance(data["dataTable"], list) or len(data["dataTable"]) == 0:
            raise ValueError("dataTable must be non-empty array")
        
        # Validate recommendations array
        if not isinstance(data["recommendations"], list) or len(data["recommendations"]) == 0:
            raise ValueError("recommendations must be non-empty array")
    
    def _load_template(self) -> str:
        """
        Load HTML template.
        
        Returns:
            str: Template HTML
            
        Raises:
            FileNotFoundError: If template not found
        """
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {self.template_path}")
        
        return self.template_path.read_text(encoding="utf-8")
    
    def _inject_data(self, template: str, title: str, data: Dict[str, Any]) -> str:
        """
        Inject data into template.
        
        Args:
            template: HTML template string
            title: Dashboard title
            data: Dashboard data
            
        Returns:
            str: Complete HTML with injected data
        """
        # Replace placeholders
        html = template.replace("{{TITLE}}", self._escape_html(title))
        html = html.replace("{{DASHBOARD_DATA}}", json.dumps(data, indent=2))
        html = html.replace("{{GENERATED_AT}}", datetime.utcnow().isoformat() + "Z")
        html = html.replace("{{VERSION}}", data["metadata"]["version"])
        
        return html
    
    def _escape_html(self, text: str) -> str:
        """
        Escape HTML special characters.
        
        Args:
            text: Text to escape
            
        Returns:
            str: Escaped text
        """
        return (text
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
            .replace("'", "&#39;"))
    
    def export_to_pdf(self, html_file: str, pdf_file: str) -> bool:
        """
        Export dashboard to PDF using Playwright.
        
        Args:
            html_file: Input HTML file path
            pdf_file: Output PDF file path
            
        Returns:
            bool: True if export successful
        """
        try:
            # Dynamic import to avoid dependency errors
            from playwright.sync_api import sync_playwright
            
            html_path = Path(html_file).resolve()
            pdf_path = Path(pdf_file)
            pdf_path.parent.mkdir(parents=True, exist_ok=True)
            
            with sync_playwright() as p:
                browser = p.chromium.launch()
                page = browser.new_page()
                page.goto(f"file://{html_path}")
                
                # Wait for visualizations to load
                page.wait_for_timeout(2000)
                
                # Export to PDF
                page.pdf(
                    path=str(pdf_path),
                    format="A4",
                    print_background=True,
                    margin={"top": "20px", "right": "20px", "bottom": "20px", "left": "20px"}
                )
                
                browser.close()
            
            return True
            
        except ImportError:
            print("Playwright not installed. Run: pip install playwright && playwright install chromium")
            return False
        except Exception as e:
            print(f"PDF export failed: {e}")
            return False
    
    def export_to_png(self, html_file: str, png_file: str) -> bool:
        """
        Export dashboard to PNG using Playwright.
        
        Args:
            html_file: Input HTML file path
            png_file: Output PNG file path
            
        Returns:
            bool: True if export successful
        """
        try:
            from playwright.sync_api import sync_playwright
            
            html_path = Path(html_file).resolve()
            png_path = Path(png_file)
            png_path.parent.mkdir(parents=True, exist_ok=True)
            
            with sync_playwright() as p:
                browser = p.chromium.launch()
                page = browser.new_page(viewport={"width": 1920, "height": 1080})
                page.goto(f"file://{html_path}")
                
                # Wait for visualizations to load
                page.wait_for_timeout(2000)
                
                # Export to PNG
                page.screenshot(
                    path=str(png_path),
                    full_page=True
                )
                
                browser.close()
            
            return True
            
        except ImportError:
            print("Playwright not installed. Run: pip install playwright && playwright install chromium")
            return False
        except Exception as e:
            print(f"PNG export failed: {e}")
            return False
    
    def export_to_pptx(self, html_file: str, pptx_file: str) -> bool:
        """
        Export dashboard to PPTX (one slide per tab).
        
        Args:
            html_file: Input HTML file path
            pptx_file: Output PPTX file path
            
        Returns:
            bool: True if export successful
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from playwright.sync_api import sync_playwright
            import tempfile
            
            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            html_path = Path(html_file).resolve()
            
            # Tab names
            tabs = ["overview", "visualizations", "diagrams", "data", "recommendations"]
            tab_titles = ["Overview", "Visualizations", "Diagrams", "Data Tables", "Recommendations"]
            
            with sync_playwright() as p:
                browser = p.chromium.launch()
                page = browser.new_page(viewport={"width": 1920, "height": 1080})
                page.goto(f"file://{html_path}")
                page.wait_for_timeout(2000)
                
                # Capture each tab
                for tab_id, tab_title in zip(tabs, tab_titles):
                    # Switch to tab
                    page.click(f"button[data-tab='{tab_id}']")
                    page.wait_for_timeout(500)
                    
                    # Take screenshot
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                        page.screenshot(path=tmp.name)
                        
                        # Add slide with image
                        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                        
                        # Add title
                        title_box = slide.shapes.add_textbox(
                            Inches(0.5), Inches(0.3),
                            Inches(9), Inches(0.5)
                        )
                        title_frame = title_box.text_frame
                        title_frame.text = tab_title
                        title_para = title_frame.paragraphs[0]
                        title_para.font.size = Pt(28)
                        title_para.font.bold = True
                        
                        # Add screenshot
                        slide.shapes.add_picture(
                            tmp.name,
                            Inches(0.5), Inches(1),
                            width=Inches(9)
                        )
                        
                        # Clean up temp file
                        Path(tmp.name).unlink()
                
                browser.close()
            
            # Save presentation
            pptx_path = Path(pptx_file)
            pptx_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(pptx_path))
            
            return True
            
        except ImportError as e:
            print(f"Missing dependencies: {e}")
            print("Run: pip install python-pptx playwright && playwright install chromium")
            return False
        except Exception as e:
            print(f"PPTX export failed: {e}")
            return False


# Example usage
if __name__ == "__main__":
    # Sample dashboard data
    sample_data = {
        "metadata": {
            "generatedAt": datetime.utcnow().isoformat() + "Z",
            "version": "3.4.0",
            "operationType": "system_alignment",
            "author": "CORTEX"
        },
        "overview": {
            "executiveSummary": "System validation completed successfully with 95% health score. All 7 layers passed validation with no critical errors. Minor warnings detected in 2 components requiring attention within next sprint.",
            "keyMetrics": [
                {
                    "label": "System Health",
                    "value": "95%",
                    "trend": "up",
                    "trendValue": "+3%",
                    "status": "healthy"
                },
                {
                    "label": "Layers Passed",
                    "value": "7/7",
                    "trend": "stable",
                    "status": "healthy"
                },
                {
                    "label": "Critical Errors",
                    "value": "0",
                    "trend": "down",
                    "trendValue": "-2",
                    "status": "healthy"
                },
                {
                    "label": "Warnings",
                    "value": "2",
                    "trend": "stable",
                    "status": "warning"
                },
                {
                    "label": "Processing Time",
                    "value": "3.2s",
                    "trend": "down",
                    "trendValue": "-0.5s",
                    "status": "healthy"
                }
            ],
            "statusIndicator": {
                "status": "healthy",
                "message": "All systems operational with no critical issues"
            }
        },
        "visualizations": {
            "forceGraph": {
                "nodes": [
                    {"id": "layer1", "group": 1, "label": "Layer 1: SKULL"},
                    {"id": "layer2", "group": 1, "label": "Layer 2: TDD"},
                    {"id": "layer3", "group": 1, "label": "Layer 3: Git"},
                    {"id": "layer4", "group": 1, "label": "Layer 4: Brain"},
                    {"id": "layer5", "group": 1, "label": "Layer 5: Agents"},
                    {"id": "layer6", "group": 1, "label": "Layer 6: Templates"},
                    {"id": "layer7", "group": 1, "label": "Layer 7: Docs"}
                ],
                "links": [
                    {"source": "layer1", "target": "layer2", "value": 1},
                    {"source": "layer2", "target": "layer3", "value": 1},
                    {"source": "layer3", "target": "layer4", "value": 1},
                    {"source": "layer4", "target": "layer5", "value": 1},
                    {"source": "layer5", "target": "layer6", "value": 1},
                    {"source": "layer6", "target": "layer7", "value": 1}
                ]
            },
            "timeSeries": {
                "labels": ["2025-11-21", "2025-11-22", "2025-11-23", "2025-11-24", "2025-11-25", "2025-11-26", "2025-11-27", "2025-11-28"],
                "datasets": [
                    {
                        "label": "System Health",
                        "data": [88, 90, 91, 92, 93, 94, 95, 95],
                        "borderColor": "rgb(75, 192, 192)",
                        "backgroundColor": "rgba(75, 192, 192, 0.2)"
                    }
                ]
            }
        },
        "diagrams": [
            {
                "title": "7-Layer Validation Flow",
                "mermaidCode": """graph TD
    A[Start Validation] --> B[Layer 1: SKULL Rules]
    B --> C[Layer 2: TDD Enforcement]
    C --> D[Layer 3: Git Isolation]
    D --> E[Layer 4: Brain Integrity]
    E --> F[Layer 5: Agent Health]
    F --> G[Layer 6: Template Compliance]
    G --> H[Layer 7: Documentation]
    H --> I[Generate Report]""",
                "type": "flowchart",
                "description": "Complete validation workflow"
            }
        ],
        "dataTable": [
            {
                "name": "Layer 1: SKULL Rules",
                "type": "Governance",
                "status": "healthy",
                "health": 100,
                "lastUpdated": "2025-11-28"
            },
            {
                "name": "Layer 2: TDD Enforcement",
                "type": "Testing",
                "status": "healthy",
                "health": 98,
                "lastUpdated": "2025-11-28"
            },
            {
                "name": "Layer 3: Git Isolation",
                "type": "Version Control",
                "status": "healthy",
                "health": 95,
                "lastUpdated": "2025-11-28"
            },
            {
                "name": "Layer 4: Brain Integrity",
                "type": "Architecture",
                "status": "warning",
                "health": 88,
                "lastUpdated": "2025-11-28"
            },
            {
                "name": "Layer 5: Agent Health",
                "type": "Operations",
                "status": "healthy",
                "health": 92,
                "lastUpdated": "2025-11-28"
            },
            {
                "name": "Layer 6: Template Compliance",
                "type": "Standards",
                "status": "warning",
                "health": 90,
                "lastUpdated": "2025-11-28"
            },
            {
                "name": "Layer 7: Documentation",
                "type": "Quality",
                "status": "healthy",
                "health": 97,
                "lastUpdated": "2025-11-28"
            }
        ],
        "recommendations": [
            {
                "priority": "medium",
                "title": "Improve Brain Integrity Score",
                "rationale": "Layer 4 health at 88% due to minor schema inconsistencies",
                "steps": [
                    "Run brain integrity check diagnostic",
                    "Review schema migration history",
                    "Apply corrective migrations if needed",
                    "Verify health improves to 95%+"
                ],
                "expectedImpact": "+7% Layer 4 health improvement",
                "estimatedEffort": "2-3 hours",
                "relatedResources": [
                    "/docs/brain-integrity-guide.md"
                ]
            },
            {
                "priority": "medium",
                "title": "Update Template Compliance",
                "rationale": "Layer 6 at 90% due to 2 templates missing response format",
                "steps": [
                    "Identify non-compliant templates",
                    "Apply response format updates",
                    "Validate against response-format.md",
                    "Re-run validation"
                ],
                "expectedImpact": "+10% Layer 6 health improvement",
                "estimatedEffort": "1-2 hours",
                "relatedResources": [
                    "/docs/response-format.md"
                ]
            },
            {
                "priority": "low",
                "title": "Schedule Next System Alignment",
                "rationale": "Proactive health monitoring recommended every 2 weeks",
                "steps": [
                    "Schedule next validation for December 12, 2025",
                    "Set calendar reminder",
                    "Prepare alignment checklist"
                ],
                "expectedImpact": "Maintain 95%+ system health",
                "estimatedEffort": "15 mins",
                "relatedResources": []
            }
        ]
    }
    
    # Generate dashboard
    generator = InteractiveDashboardGenerator()
    output_file = "cortex-brain/documents/reports/sample-dashboard.html"
    
    print("Generating dashboard...")
    success = generator.generate_dashboard(
        title="CORTEX System Alignment Dashboard",
        data=sample_data,
        output_file=output_file
    )
    
    if success:
        print(f"✅ Dashboard generated: {output_file}")
        
        # Export examples
        print("\nExporting to PDF...")
        generator.export_to_pdf(output_file, "cortex-brain/documents/reports/sample-dashboard.pdf")
        
        print("Exporting to PNG...")
        generator.export_to_png(output_file, "cortex-brain/documents/reports/sample-dashboard.png")
        
        print("Exporting to PPTX...")
        generator.export_to_pptx(output_file, "cortex-brain/documents/reports/sample-dashboard.pptx")
    else:
        print("❌ Dashboard generation failed")
