"""DocumentReader — Phase 144-b.

Extracts text from Office documents and PDFs using lazy-loaded libraries.
Gracefully degrades if Office libraries (python-docx, openpyxl,
python-pptx, pypdf) are not installed.

Source: GitHub Issue #17 — FB-2026-03-09-074435-001
CORE: CORE-008, CORE-011, CORE-012
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, List, Optional

logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────────────────────────
# Data classes
# ─────────────────────────────────────────────────────────────────────────────

@dataclass
class DocumentContent:
    """Extracted content from a document.

    Attributes:
        path: Source file path.
        format: File format (e.g. 'docx', 'xlsx', 'pdf').
        title: Inferred document title (first heading or filename stem).
        text: Full extracted text content.
        sections: Ordered list of (heading, body) tuples.
        page_count: Number of pages or sheets (0 if unknown).
        error: Non-empty if extraction failed; text will be empty.
    """

    path: Path
    format: str
    title: str = ""
    text: str = ""
    sections: List[Dict[str, str]] = field(default_factory=list)
    page_count: int = 0
    error: str = ""


# ─────────────────────────────────────────────────────────────────────────────
# Reader
# ─────────────────────────────────────────────────────────────────────────────

class DocumentReader:
    """Reads Office and PDF documents into DocumentContent.

    All Office library imports are deferred to method-call time so the
    module remains importable even when libraries are absent (graceful
    degradation — ``error`` field populated instead of raising).

    Supported formats:
        - ``.docx`` / ``.doc`` — via python-docx
        - ``.xlsx`` / ``.xls`` — via openpyxl
        - ``.pptx`` / ``.ppt`` — via python-pptx
        - ``.pdf``             — via pypdf
        - ``.md`` / ``.txt`` / ``.rst`` / ``.yaml`` / ``.yml`` — built-in

    Usage::

        reader = DocumentReader()
        content = reader.read(Path("report.docx"))
    """

    def read(self, path: Path) -> DocumentContent:
        """Dispatch to the appropriate reader based on file extension.

        Args:
            path: Path to the document file.

        Returns:
            DocumentContent with extracted text and metadata.
        """
        suffix = path.suffix.lower()
        dispatch = {
            ".docx": self.read_docx,
            ".doc": self.read_docx,
            ".xlsx": self.read_xlsx,
            ".xls": self.read_xlsx,
            ".pptx": self.read_pptx,
            ".ppt": self.read_pptx,
            ".pdf": self.read_pdf,
            ".md": self._read_text,
            ".txt": self._read_text,
            ".rst": self._read_text,
            ".yaml": self._read_text,
            ".yml": self._read_text,
        }
        handler = dispatch.get(suffix, self._read_unsupported)
        return handler(path)

    # ── Format-specific readers ───────────────────────────────────────────

    def read_docx(self, path: Path) -> DocumentContent:
        """Extract paragraphs from a Word document via python-docx.

        Args:
            path: Path to .docx file.

        Returns:
            DocumentContent with each paragraph as text.
        """
        fmt = path.suffix.lstrip(".").lower() or "docx"
        content = DocumentContent(path=path, format=fmt)
        try:
            import docx  # python-docx exposes 'docx' module
            doc = docx.Document(str(path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            content.text = "\n".join(paragraphs)
            content.title = paragraphs[0] if paragraphs else path.stem
            content.page_count = len(doc.paragraphs)
            current_section: Dict[str, str] = {}
            for para in doc.paragraphs:
                if para.style.name.startswith("Heading"):
                    if current_section:
                        content.sections.append(current_section)
                    current_section = {"heading": para.text, "body": ""}
                elif current_section:
                    current_section["body"] += para.text + "\n"
            if current_section:
                content.sections.append(current_section)
        except ImportError:
            content.error = "python-docx not installed — install with: pip install python-docx"
            logger.warning("DocumentReader.read_docx: %s", content.error)
        except Exception as exc:  # noqa: BLE001
            content.error = f"Failed to read Word document: {exc}"
            logger.warning("DocumentReader.read_docx: %s", content.error)
        return content

    def read_xlsx(self, path: Path) -> DocumentContent:
        """Extract cell values from an Excel workbook via openpyxl.

        Args:
            path: Path to .xlsx file.

        Returns:
            DocumentContent with tab-separated cell values per row.
        """
        fmt = path.suffix.lstrip(".").lower() or "xlsx"
        content = DocumentContent(path=path, format=fmt, title=path.stem)
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            rows_text: List[str] = []
            for sheet in wb.worksheets:
                content.sections.append({"heading": sheet.title, "body": ""})
                sheet_rows: List[str] = []
                for row in sheet.iter_rows(values_only=True):
                    cells = "\t".join(str(c) if c is not None else "" for c in row)
                    if cells.strip():
                        sheet_rows.append(cells)
                body = "\n".join(sheet_rows)
                content.sections[-1]["body"] = body
                rows_text.append(f"[{sheet.title}]\n{body}")
            content.text = "\n\n".join(rows_text)
            content.page_count = len(wb.worksheets)
            wb.close()
        except ImportError:
            content.error = "openpyxl not installed — install with: pip install openpyxl"
            logger.warning("DocumentReader.read_xlsx: %s", content.error)
        except Exception as exc:  # noqa: BLE001
            content.error = f"Failed to read Excel workbook: {exc}"
            logger.warning("DocumentReader.read_xlsx: %s", content.error)
        return content

    def read_pptx(self, path: Path) -> DocumentContent:
        """Extract slide text from a PowerPoint presentation via python-pptx.

        Args:
            path: Path to .pptx file.

        Returns:
            DocumentContent with each slide's text as a section.
        """
        fmt = path.suffix.lstrip(".").lower() or "pptx"
        content = DocumentContent(path=path, format=fmt, title=path.stem)
        try:
            from pptx import Presentation  # python-pptx
            prs = Presentation(str(path))
            slide_texts: List[str] = []
            for i, slide in enumerate(prs.slides, start=1):
                slide_text = " ".join(
                    shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()
                )
                content.sections.append({"heading": f"Slide {i}", "body": slide_text})
                slide_texts.append(slide_text)
            content.text = "\n".join(slide_texts)
            content.page_count = len(prs.slides)
            if content.sections:
                content.title = content.sections[0].get("body", path.stem).split("\n")[0][:80]
        except ImportError:
            content.error = "python-pptx not installed — install with: pip install python-pptx"
            logger.warning("DocumentReader.read_pptx: %s", content.error)
        except Exception as exc:  # noqa: BLE001
            content.error = f"Failed to read PowerPoint presentation: {exc}"
            logger.warning("DocumentReader.read_pptx: %s", content.error)
        return content

    def read_pdf(self, path: Path) -> DocumentContent:
        """Extract page text from a PDF via pypdf.

        Args:
            path: Path to .pdf file.

        Returns:
            DocumentContent with each page as a section.
        """
        content = DocumentContent(path=path, format="pdf", title=path.stem)
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            pages: List[str] = []
            for i, page in enumerate(reader.pages, start=1):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    content.sections.append({"heading": f"Page {i}", "body": page_text})
                    pages.append(page_text)
            content.text = "\n".join(pages)
            content.page_count = len(reader.pages)
            if pages:
                content.title = pages[0].split("\n")[0][:80]
        except ImportError:
            content.error = "pypdf not installed — install with: pip install pypdf"
            logger.warning("DocumentReader.read_pdf: %s", content.error)
        except Exception as exc:  # noqa: BLE001
            content.error = f"Failed to read PDF: {exc}"
            logger.warning("DocumentReader.read_pdf: %s", content.error)
        return content

    # ── Internal helpers ─────────────────────────────────────────────────

    def _read_text(self, path: Path) -> DocumentContent:
        """Read plain text / Markdown / YAML file using built-in open()."""
        fmt = path.suffix.lstrip(".").lower() or "txt"
        content = DocumentContent(path=path, format=fmt, title=path.stem)
        try:
            raw = path.read_text(encoding="utf-8", errors="replace")
            content.text = raw
            lines = raw.splitlines()
            # Extract first heading as title
            for line in lines:
                stripped = line.strip().lstrip("#").strip()
                if stripped:
                    content.title = stripped[:80]
                    break
        except Exception as exc:  # noqa: BLE001
            content.error = f"Failed to read text file: {exc}"
        return content

    def _read_unsupported(self, path: Path) -> DocumentContent:
        """Return error content for unsupported file types."""
        return DocumentContent(
            path=path,
            format=path.suffix.lstrip(".").lower() or "unknown",
            error=f"Unsupported file format: '{path.suffix}'",
        )
